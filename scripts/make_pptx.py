"""
Sinh file PPT báo cáo cuối khóa từ nội dung DOCX.
Input: reports/2582003023_BaoCaoCuoiKy_LTCKHDL_NguyenMinhHuy.docx + figures/
Output: reports/2582003023_BaoCaoCuoiKy_Presentation.pptx
"""

import os, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from copy import deepcopy

BASE = Path(__file__).parent.parent.resolve()
REPORTS = BASE / "reports"
FIG_DIR = REPORTS / "figures"
EXTRACT_DIR = REPORTS / "figures_extracted"
OUT = REPORTS / "2582003023_BaoCaoCuoiKy_Presentation.pptx"

# ── Màu sắc ──────────────────────────────────────────────────────────────
DARK = RGBColor(0x1F, 0x36, 0x4D)      # navy
ACCENT = RGBColor(0xC1, 0x21, 0x1B)    # đỏ
LIGHT = RGBColor(0xEA, 0xEC, 0xF0)     # xám nhạt
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x2B, 0x2B, 0x2B)
SUBTLE = RGBColor(0x66, 0x6C, 0x7B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

# ── Slides 16:9 ──────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, color, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font="Calibri", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_runs(slide, x, y, w, h, runs, *,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs: list of (text, size, bold, color, italic?)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(runs):
        text = item[0]
        size = item[1] if len(item) > 1 else 16
        bold = item[2] if len(item) > 2 else False
        color = item[3] if len(item) > 3 else TEXT
        italic = item[4] if len(item) > 4 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def add_bullet_box(slide, x, y, w, h, items, *,
                   size=18, color=TEXT, bullet_color=ACCENT, line_spacing=1.15):
    """items: list of strings, or list of (str, sub_size)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, sz = item[0], item[1]
        else:
            text, sz = item, size
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(6)
        # bullet char
        rb = p.add_run()
        rb.text = "▸ "
        rb.font.size = Pt(sz)
        rb.font.bold = True
        rb.font.color.rgb = bullet_color
        rb.font.name = "Calibri"
        rt = p.add_run()
        rt.text = text
        rt.font.size = Pt(sz)
        rt.font.color.rgb = color
        rt.font.name = "Calibri"
    return tb


def add_image(slide, path, x, y, w=None, h=None):
    if not Path(path).exists():
        return None
    if w and h:
        return slide.shapes.add_picture(str(path), x, y, width=w, height=h)
    if w:
        return slide.shapes.add_picture(str(path), x, y, width=w)
    return slide.shapes.add_picture(str(path), x, y, height=h)


def slide_header(slide, title, subtitle=None, page=None, total=None):
    # Top accent bar
    add_rect(slide, 0, 0, SW, Inches(0.06), DARK)
    # Title
    add_text(slide, Inches(0.5), Inches(0.20), SW - Inches(1), Inches(0.7),
             title, size=28, bold=True, color=DARK)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.85), SW - Inches(1), Inches(0.4),
                 subtitle, size=14, color=SUBTLE, italic=True)
    # Footer line
    add_rect(slide, 0, SH - Inches(0.04), SW, Inches(0.04), DARK)
    # Footer
    add_text(slide, Inches(0.5), SH - Inches(0.4), Inches(7), Inches(0.3),
             "Nguyễn Minh Huy - 2582003023 - LTCKHDL", size=10, color=SUBTLE)
    if page and total:
        add_text(slide, SW - Inches(2.0), SH - Inches(0.4), Inches(1.5), Inches(0.3),
                 f"Slide {page} / {total}", size=10, color=SUBTLE, align=PP_ALIGN.RIGHT)


def slide_blank():
    return prs.slides.add_slide(BLANK)


def slide_title_big(title, subtitle=""):
    s = slide_blank()
    # Big navy background block
    add_rect(s, 0, 0, SW, SH, DARK)
    # Accent strip
    add_rect(s, 0, Inches(6.6), SW, Inches(0.9), ACCENT)
    # Title
    add_text(s, Inches(0.8), Inches(2.5), SW - Inches(1.6), Inches(1.4),
             title, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(s, Inches(0.8), Inches(4.0), SW - Inches(1.6), Inches(0.8),
                 subtitle, size=22, color=WHITE, italic=True, align=PP_ALIGN.CENTER)
    # Bottom credit
    add_text(s, Inches(0.8), Inches(6.7), SW - Inches(1.6), Inches(0.7),
             "Nguyễn Minh Huy - 2582003023 - LTCKHDL", size=14, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


def slide_section(title):
    s = slide_blank()
    add_rect(s, 0, 0, SW, SH, DARK)
    # Accent stripe on left
    add_rect(s, 0, Inches(3), Inches(0.4), Inches(1.5), ACCENT)
    add_text(s, Inches(0.8), Inches(2.5), SW - Inches(1.6), Inches(1.5),
             title, size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.MIDDLE)
    return s


# ─────────────────────────────────────────────────────────────────────────
# SLIDE 1: BÌA
# ─────────────────────────────────────────────────────────────────────────
s = slide_title_big(
    "PHÂN TÍCH VÀ GỢI Ý BẤT ĐỘNG SẢN\nCĂN HỘ / CHUNG CƯ TẠI TP. HỒ CHÍ MINH",
    "Báo cáo cuối kỳ - Môn Lập trình cho Khoa học Dữ liệu"
)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 2: GIỚI THIỆU
# ─────────────────────────────────────────────────────────────────────────
total_slides = 30  # sẽ điều chỉnh sau

s = slide_blank()
slide_header(s, "GIỚI THIỆU SINH VIÊN & ĐỀ TÀI", page=2, total=total_slides)
# Two columns
add_rect(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.3), LIGHT)
add_rect(s, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.3), LIGHT)
# Left: Student info
add_text(s, Inches(0.7), Inches(1.5), Inches(5.6), Inches(0.5),
         "THÔNG TIN SINH VIÊN", size=18, bold=True, color=DARK)
add_bullet_box(s, Inches(0.7), Inches(2.1), Inches(5.6), Inches(4.5), [
    "Học viên cao học: Nguyễn Minh Huy",
    "MSSV: 2582003023",
    "Lớp: 33CNTT21-PH",
    "Ngành: Công Nghệ Thông Tin",
    "GVHD: PGS.TS Nguyễn Duy Hàm",
    "Trường: Đại học Thủy Lợi",
    "TP.HCM, tháng 8 năm 2026",
], size=18)
# Right: Topic
add_text(s, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.5),
         "ĐỀ TÀI", size=18, bold=True, color=DARK)
add_text(s, Inches(7.0), Inches(2.1), Inches(5.6), Inches(0.7),
         "Chuyên đề 3", size=22, bold=True, color=ACCENT)
add_text(s, Inches(7.0), Inches(2.7), Inches(5.6), Inches(1.3),
         "Phân tích và gợi ý BĐS căn hộ/chung cư TP.HCM",
         size=16, bold=True, color=TEXT)
add_bullet_box(s, Inches(7.0), Inches(3.7), Inches(5.6), Inches(2.5), [
    "Dataset: 1.799 tin (chotot.com, T7-T8/2026)",
    "Mục tiêu: Chuẩn hóa → EDA → Dự đoán → Phân cụm → Gợi ý",
    "6 quận TP.HCM: Thủ Đức, Bình Thạnh, Q7, Gò Vấp, Q12, Bình Tân",
], size=14, bullet_color=DARK)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 3: MỤC LỤC
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "MỤC LỤC BÁO CÁO", page=3, total=total_slides)
toc_items = [
    ("LỜI MỞ ĐẦU", "Lý do, mục đích, phạm vi nghiên cứu"),
    ("CHƯƠNG 1", "Giới thiệu bài toán"),
    ("CHƯƠNG 2", "Tổng quan lý thuyết nền tảng"),
    ("CHƯƠNG 3", "Triển khai thuật toán"),
    ("CHƯƠNG 4", "Thực nghiệm và kết quả"),
    ("KẾT LUẬN", "Tổng kết và hướng phát triển"),
]
y0 = Inches(1.6)
for i, (sec, desc) in enumerate(toc_items):
    add_rect(s, Inches(0.6), y0 + Inches(i*0.85), Inches(0.15), Inches(0.7),
             ACCENT if i % 2 == 0 else DARK)
    add_text(s, Inches(1.0), y0 + Inches(i*0.85), Inches(4.5), Inches(0.7),
             sec, size=22, bold=True, color=DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.5), y0 + Inches(i*0.85), Inches(7.5), Inches(0.7),
             desc, size=18, color=TEXT,
             anchor=MSO_ANCHOR.MIDDLE)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 4: LỜI MỞ ĐẦU - Lý do
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "LỜI MỞ ĐẦU", subtitle="Lý do lựa chọn đề tài",
             page=4, total=total_slides)
add_bullet_box(s, Inches(0.7), Inches(1.4), Inches(12), Inches(5.5), [
    ("Thị trường căn hộ TP.HCM duy trì mức giao dịch sôi động, hàng nghìn tin/ngày trên chotot, batdongsan, alonhadat.", 18),
    ("Dữ liệu thô chưa chuẩn hóa: giá viết dạng text (\"3,19 tỷ\"), hướng nhà mã 1–8, nội thất mã 1–4, pháp lý mã 1,2,4,5,6.", 18),
    ("Outlier rõ ràng: căn 1323m² (gấp 16× trung bình), căn có 11 phòng ngủ (không hợp lý với căn hộ).", 18),
    ("Missing cao: direction 79%, balcony 73%, furnishing 34% → cần imputation hợp lý.", 18),
    ("Cần giải pháp thực tiễn để chuẩn hóa dữ liệu, dự đoán giá, phân cụm thị trường và gợi ý sản phẩm phù hợp.", 18),
], size=18)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 5: Mục đích nghiên cứu
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "MỤC ĐÍCH NGHIÊN CỨU", page=5, total=total_slides)
add_bullet_box(s, Inches(0.7), Inches(1.4), Inches(12), Inches(5.5), [
    ("Làm sạch & chuẩn hóa bộ dữ liệu tin đăng căn hộ tại 6 quận TP.HCM từ nguồn snapshot.", 18),
    ("Nhận diện và xử lý giá trị ngoại lệ (outlier) ở diện tích, số phòng, giá bán.", 18),
    ("Huấn luyện các thuật toán học máy (Linear Regression, Random Forest, Gradient Boosting) để dự đoán giá/m².", 18),
    ("Phân chia phân khúc thị trường bằng K-Means, tự động chọn K bằng chỉ số Silhouette.", 18),
    ("Phát triển hệ gợi ý đề xuất top-5 tin đăng sát nhất với hồ sơ nhu cầu người dùng.", 18),
    ("Trực quan hóa dữ liệu bằng các biểu đồ có nhận xét phân tích chuyên sâu.", 18),
], size=18)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 6: Phương pháp & Phạm vi
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "PHƯƠNG PHÁP & PHẠM VI NGHIÊN CỨU", page=6, total=total_slides)
add_rect(s, Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5), LIGHT)
add_rect(s, Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5), LIGHT)
add_text(s, Inches(0.7), Inches(1.5), Inches(5.6), Inches(0.5),
         "PHƯƠNG PHÁP", size=18, bold=True, color=DARK)
add_bullet_box(s, Inches(0.7), Inches(2.1), Inches(5.6), Inches(4.5), [
    "Quy trình CRISP-DM 9 bước",
    "Python + pandas, numpy, scikit-learn",
    "5-fold Cross-Validation trên train",
    "Đánh giá độc lập trên tập test 20%",
    "Log1p transform cho biến target",
    "Silhouette auto-pick K cho K-Means",
], size=16)
add_text(s, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.5),
         "PHẠM VI", size=18, bold=True, color=DARK)
add_bullet_box(s, Inches(7.0), Inches(2.1), Inches(5.6), Inches(4.5), [
    "1.799 tin căn hộ từ chotot.com (T7-T8/2026)",
    "6 quận TP.HCM: Thủ Đức, Bình Thạnh, Q7, Gò Vấp, Q12, Bình Tân",
    "Sau làm sạch: 1.773 dòng (dropna price)",
    "Snapshot 1 thời điểm, không phân tích xu hướng",
    "Bổ sung nguồn phụ: 91 dòng tiện ích (trường, bệnh viện, công viên)",
], size=16)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 7: SECTION - CHƯƠNG 1
# ─────────────────────────────────────────────────────────────────────────
s = slide_section("CHƯƠNG 1\nGiới thiệu bài toán")

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 8: 6 câu hỏi nghiên cứu
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "6 CÂU HỎI NGHIÊN CỨU", page=8, total=total_slides)
qs = [
    "Trong 6 quận khảo sát, khu vực nào có giá/m² cao nhất và thấp nhất?",
    "Diện tích và số phòng ngủ ảnh hưởng thế nào đến giá/m²?",
    "Hướng nhà, ban công, nội thất, pháp lý liên hệ thống kê ra sao với giá?",
    "Hiệu suất mô hình (MAE, RMSE, R²) đạt mức nào trên tập test?",
    "Thị trường căn hộ TP.HCM có thể phân thành bao nhiêu phân khúc tối ưu?",
    "Hệ thống gợi ý có đề xuất đúng top-5 tin đăng theo hồ sơ nhu cầu?",
]
for i, q in enumerate(qs):
    y = Inches(1.3) + Inches(i * 0.95)
    # Question number circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y, Inches(0.6), Inches(0.6))
    circ.fill.solid()
    circ.fill.fore_color.rgb = ACCENT
    circ.line.fill.background()
    add_text(s, Inches(0.6), y, Inches(0.6), Inches(0.6),
             f"C{i+1}", size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Question text
    add_rect(s, Inches(1.4), y, Inches(11.3), Inches(0.75), LIGHT)
    add_text(s, Inches(1.6), y, Inches(11.1), Inches(0.75),
             q, size=15, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 9: Từ điển dữ liệu (Table 2)
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "TỪ ĐIỂN DỮ LIỆU", subtitle="16 thuộc tính chính",
             page=9, total=total_slides)
# Build table
rows = [
    ("listing_id", "int", "0.0%", "Khóa chính", "Mã tin duy nhất"),
    ("total_price", "int", "0.0%", "Target phụ", "Tổng giá (VND)"),
    ("area_m2", "float", "0.3%", "Feature chính", "Diện tích (m²)"),
    ("price_per_m2", "float", "0.3%", "Target chính", "Giá/m² (VND), tính lại từ total/area"),
    ("bedrooms", "float", "0.2%", "Feature chính", "Số phòng ngủ"),
    ("bathrooms", "float", "15.6%", "Feature phụ", "Số phòng vệ sinh"),
    ("direction", "1–8", "78.9%", "Feature", "Mã hướng nhà"),
    ("balcony_dir", "1–8", "73.2%", "Feature", "Mã hướng ban công"),
    ("furnishing", "1–4", "34.4%", "Feature", "Mã nội thất"),
    ("legal_status", "1,2,4,5,6", "20.0%", "Feature", "Mã pháp lý"),
    ("apt_type", "int", "0.0%", "Feature", "Loại căn hộ"),
    ("image_count", "int", "0.0%", "Feature", "Số ảnh - chất lượng tin"),
    ("project_name", "str", "20.9%", "Categorical", "Tên dự án → OHE"),
    ("district_clean", "str", "0.0%", "Categorical", "Quận (6 quận)"),
    ("amenity_score", "float", "0.0%", "Feature", "Tiện ích phường (nguồn phụ)"),
    ("cluster", "int 0–3", "0.0%", "Phân cụm", "K-Means K=4"),
]
n_rows = 1 + len(rows)
table = s.shapes.add_table(n_rows, 5, Inches(0.5), Inches(1.3),
                            Inches(12.3), Inches(5.7)).table
headers = ["Thuộc tính", "Kiểu", "Missing", "Vai trò", "Mô tả"]
col_widths = [Inches(2.0), Inches(1.4), Inches(1.2), Inches(1.8), Inches(5.9)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w
# Header
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True
            r.font.color.rgb = WHITE
            r.font.size = Pt(13)
            r.font.name = "Calibri"
# Rows
for ri, row in enumerate(rows):
    bg = LIGHT if ri % 2 == 0 else WHITE
    for ci, val in enumerate(row):
        cell = table.cell(ri+1, ci)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT if ci > 0 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(11)
                r.font.color.rgb = TEXT
                r.font.bold = (ci == 0)
                r.font.name = "Calibri"

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 10: Quy trình thu thập & làm sạch
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "QUY TRÌNH THU THẬP & LÀM SẠCH", page=10, total=total_slides)
# Flow steps
steps = [
    ("1", "Crawl", "API chotot", "→ 1799"),
    ("2", "Làm sạch", "5 bước", "→ 1773"),
    ("3", "EDA", "10 biểu đồ", "→ Nhận xét"),
    ("4", "ML", "4 mô hình", "→ R²=0.169"),
    ("5", "Phân cụm", "K-Means", "→ K=4"),
    ("6", "Gợi ý", "Hybrid filter", "→ Top-5"),
]
y0 = Inches(2.2)
for i, (n, t1, t2, t3) in enumerate(steps):
    x = Inches(0.5) + Inches(i * 2.2)
    add_rect(s, x, y0, Inches(2.0), Inches(2.5), LIGHT)
    add_rect(s, x, y0, Inches(2.0), Inches(0.4), DARK)
    add_text(s, x, y0, Inches(2.0), Inches(0.4),
             n, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, y0 + Inches(0.5), Inches(2.0), Inches(0.5),
             t1, size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(1.0), Inches(2.0), Inches(0.4),
             t2, size=12, color=SUBTLE, align=PP_ALIGN.CENTER)
    add_text(s, x, y0 + Inches(1.5), Inches(2.0), Inches(0.4),
             t3, size=12, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

# Cleaning steps
add_text(s, Inches(0.5), Inches(5.2), Inches(12), Inches(0.5),
         "5 bước làm sạch:", size=16, bold=True, color=DARK)
add_bullet_box(s, Inches(0.7), Inches(5.7), Inches(12), Inches(1.3), [
    "(1) Parse giá text → số   (2) Chuẩn hóa district   (3) Decode direction 1–8   (4) Decode furnishing + legal   (5) Lọc outlier (3 area + 17 bedrooms)",
], size=14)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 11: Hình 1 - Listings theo quận
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "PHÂN BỐ TIN ĐĂNG THEO 6 QUẬN", page=11, total=total_slides)
img = add_image(s, FIG_DIR / "fig05_listings_by_district.png",
                Inches(2.5), Inches(1.4), w=Inches(8))
if img is None:
    add_text(s, Inches(4), Inches(3), Inches(5), Inches(1),
             "[Hình không tìm thấy]", size=20, color=ACCENT,
             align=PP_ALIGN.CENTER)
# Caption + nhận xét
add_text(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.5),
         "Hình 1: Phân bố số tin đăng theo 6 quận TP.HCM (1.773 dòng sau làm sạch).",
         size=13, italic=True, color=SUBTLE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(1.0),
         "Nhận xét: TP. Thủ Đức dẫn đầu với 593 tin (33.4%), kế tiếp là Quận 7 (352), Bình Tân (302), Bình Thạnh (212), Quận 12 (173), Gò Vấp (141).",
         size=14, color=TEXT)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 12: SECTION - CHƯƠNG 2
# ─────────────────────────────────────────────────────────────────────────
s = slide_section("CHƯƠNG 2\nTổng quan lý thuyết nền tảng")

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 13: Lý thuyết các thuật toán (gộp)
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "CƠ SỞ LÝ THUYẾT CÁC THUẬT TOÁN", page=13, total=total_slides)
theories = [
    ("Linear Regression", "Tìm hàm tuyến tính f(x)=w·x+b tối thiểu hóa tổng bình phương sai số. Áp dụng log1p transform lên target để giảm skewness, sau đó inverse bằng expm1."),
    ("Random Forest", "Ensemble nhiều cây quyết định (200 cây), mỗi cây fit trên bootstrap sample + random features. Cấu hình: n_estimators=200, min_samples_leaf=2. Tự động nắm bắt tương tác phi tuyến."),
    ("Gradient Boosting", "Xây dựng cây tuần tự, mỗi cây sửa sai số của cây trước (stage-wise). Cấu hình: n_estimators=200, max_depth=4, lr=0.05. CV R² cao nhưng dễ overfit."),
    ("K-Means", "Phân cụm unsupervised, tối ưu WCSS. Pipeline: StandardScaler → KMeans(n_init=10). Chọn K tự động bằng Silhouette s = (b-a)/max(a,b) ∈ [-1,1]."),
    ("Hybrid Recommender", "Lọc cứng (budget ±20%, bedrooms ±1, district) + chấm điểm có trọng số (price + area + cluster_bonus + amenity_bonus)."),
]
y0 = Inches(1.3)
for i, (name, desc) in enumerate(theories):
    y = y0 + Inches(i * 1.05)
    add_rect(s, Inches(0.5), y, Inches(2.8), Inches(0.95), DARK)
    add_text(s, Inches(0.5), y, Inches(2.8), Inches(0.95),
             name, size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(3.4), y, Inches(9.4), Inches(0.95), LIGHT)
    add_text(s, Inches(3.6), y, Inches(9.1), Inches(0.95),
             desc, size=12, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 14: SECTION - CHƯƠNG 3
# ─────────────────────────────────────────────────────────────────────────
s = slide_section("CHƯƠNG 3\nTriển khai thuật toán")

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 15: Kiến trúc tổng thể
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "KIẾN TRÚC TỔNG THỂ HỆ THỐNG", page=15, total=total_slides)
# 4 boxes
stages = [
    ("Thu thập & Làm sạch", "xlsx → cleaner.py\n→ outlier filter\n→ listings_clean.csv", "1799 → 1779"),
    ("Bổ sung & Transform", "merge amenities\n+ ColumnTransformer\n(47 features)", "1773 dòng"),
    ("Mô hình học máy", "Dummy/Linear/RF/GBR\n+ 5-fold CV\n+ K-Means K=4", "R²=0.169"),
    ("Gợi ý Hybrid", "Hybrid filter\n+ Scoring\n→ Top-5", "3 profiles"),
]
y0 = Inches(2.0)
for i, (title, desc, badge) in enumerate(stages):
    x = Inches(0.5) + Inches(i * 3.2)
    add_rect(s, x, y0, Inches(3.0), Inches(3.5), LIGHT)
    add_rect(s, x, y0, Inches(3.0), Inches(0.5), DARK)
    add_text(s, x, y0, Inches(3.0), Inches(0.5),
             title, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.15), y0 + Inches(0.7), Inches(2.7), Inches(2.0),
             desc, size=13, color=TEXT, align=PP_ALIGN.CENTER)
    # Badge
    add_rect(s, x + Inches(0.5), y0 + Inches(2.7), Inches(2.0), Inches(0.5), ACCENT)
    add_text(s, x + Inches(0.5), y0 + Inches(2.7), Inches(2.0), Inches(0.5),
             badge, size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.6),
         "Mã nguồn Python thuần, kiến trúc module: cleaner, data_manager, features, predictor, segmenter, recommender, pipeline (CLI end-to-end)",
         size=12, italic=True, color=SUBTLE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 16: Feature Pipeline (ColumnTransformer)
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "FEATURE PIPELINE - ColumnTransformer", page=16, total=total_slides)
# Numeric branch
add_rect(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(2.8), LIGHT)
add_text(s, Inches(0.7), Inches(1.5), Inches(5.6), Inches(0.5),
         "NHÁNH NUMERIC (9 features)", size=15, bold=True, color=DARK)
add_text(s, Inches(0.7), Inches(2.0), Inches(5.6), Inches(2.0),
         "SimpleImputer(strategy='median')\n   ↓\nStandardScaler()\n   ↓\nOutput: 9 features chuẩn hóa",
         size=13, color=TEXT, font="Consolas")
add_text(s, Inches(0.7), Inches(3.8), Inches(5.6), Inches(0.4),
         "Chọn median thay vì mean: ít bị ảnh hưởng bởi outlier",
         size=11, italic=True, color=SUBTLE)

# Categorical branch
add_rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(2.8), LIGHT)
add_text(s, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.5),
         "NHÁNH CATEGORICAL (3 features)", size=15, bold=True, color=DARK)
add_text(s, Inches(7.0), Inches(2.0), Inches(5.6), Inches(2.0),
         "SimpleImputer(strategy='constant', 'missing')\n   ↓\nOneHotEncoder(handle_unknown='ignore',\n                min_frequency=10)\n   ↓\nOutput: 38 dummy features",
         size=13, color=TEXT, font="Consolas")
add_text(s, Inches(7.0), Inches(3.8), Inches(5.6), Inches(0.4),
         "min_frequency=10 gộp category hiếm thành 1 dummy",
         size=11, italic=True, color=SUBTLE)

# Total
add_rect(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.8), DARK)
add_text(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.8),
         "TỔNG: 9 numeric + 38 OHE categorical = 47 features (fit trên train, transform độc lập test)",
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.5), Inches(5.6), Inches(12), Inches(1.3),
         "Tại sao ColumnTransformer?\n• Pipeline sklearn chuẩn → reproduce được trên dữ liệu mới\n• Fit-transform trên train → không có data leakage\n• Đảm bảo numeric/categorical được xử lý đúng cách song song",
         size=13, color=TEXT)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 17: Module PricePredictor
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "MODULE PRICE PREDICTOR - 4 MÔ HÌNH", page=17, total=total_slides)
add_text(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
         "Cấu trúc thống nhất: fit(X, y) → predict(X) → evaluate(X_test, y_test)",
         size=14, italic=True, color=SUBTLE)
# 4 models in 2x2 grid
models = [
    ("Dummy Regressor", "Baseline", "Strategy: median\nKhông học gì từ features", ACCENT),
    ("Linear Regression", "Hồi quy tuyến tính", "log1p(y) trước fit\nexpm1() sau predict", DARK),
    ("Random Forest ★", "Ensemble cây", "200 cây\nmin_samples_leaf=2\nGeneralize tốt nhất", GREEN),
    ("Gradient Boosting", "Boosting tuần tự", "200 stages\nmax_depth=4\nlr=0.05", ORANGE),
]
positions = [(0.5, 1.9), (6.85, 1.9), (0.5, 4.4), (6.85, 4.4)]
for (x, y), (name, sub, desc, color) in zip(positions, models):
    add_rect(s, Inches(x), Inches(y), Inches(6.0), Inches(2.3), LIGHT)
    add_rect(s, Inches(x), Inches(y), Inches(0.3), Inches(2.3), color)
    add_text(s, Inches(x + 0.5), Inches(y + 0.1), Inches(5.5), Inches(0.5),
             name, size=17, bold=True, color=DARK)
    add_text(s, Inches(x + 0.5), Inches(y + 0.6), Inches(5.5), Inches(0.4),
             sub, size=12, italic=True, color=SUBTLE)
    add_text(s, Inches(x + 0.5), Inches(y + 1.1), Inches(5.5), Inches(1.1),
             desc, size=12, color=TEXT, font="Consolas")

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 18: Hình 4 - Heatmap correlation
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "EDA - HEATMAP TƯƠNG QUAN", page=18, total=total_slides)
img = add_image(s, FIG_DIR / "fig06_correlation_heatmap.png",
                Inches(2.5), Inches(1.3), w=Inches(8))
if img is None:
    add_text(s, Inches(4), Inches(3), Inches(5), Inches(1),
             "[Hình không tìm thấy]", size=20, color=ACCENT)
add_text(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
         "Hình 2: Heatmap tương quan Pearson giữa các biến số.",
         size=13, italic=True, color=SUBTLE, align=PP_ALIGN.CENTER)
add_bullet_box(s, Inches(0.7), Inches(6.0), Inches(12), Inches(1.2), [
    "Tương quan giữa price_per_m2 và các biến khác yếu (<0.3) → khó dự đoán chỉ từ numeric features.",
    "bedrooms ↔ area_m2 = 0.55 (tương quan dự kiến); total_price ↔ area_m2 = 0.70.",
], size=13)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 19: SECTION - CHƯƠNG 4
# ─────────────────────────────────────────────────────────────────────────
s = slide_section("CHƯƠNG 4\nThực nghiệm và kết quả")

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 20: Trả lời C1 - Quận giá cao nhất
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "C1: QUẬN NÀO CÓ GIÁ/M² CAO NHẤT?", page=20, total=total_slides)
# Bar chart of medians
import pandas as pd
df = pd.read_csv(BASE / "data/processed/listings_with_amenities.csv").dropna(subset=["price_per_m2"])
medians = df.groupby("district_clean")["price_per_m2"].median().sort_values(ascending=False)
# Plot as bar chart on slide
chart_data = [(dist, m/1e6) for dist, m in medians.items()]
max_m = max(m for _, m in chart_data)
y0 = Inches(1.3)
bar_w = Inches(1.5)
gap = Inches(0.3)
total_w = len(chart_data) * (bar_w.emu + gap.emu)
start_x = (SW.emu - total_w) / 2
chart_h = Inches(4.0)
for i, (dist, m) in enumerate(chart_data):
    x = Emu(start_x) + i * (bar_w + gap)
    h = Emu(int(chart_h.emu * m / max_m))
    color = ACCENT if i == 0 else (DARK if i == 1 else GREEN if i == len(chart_data)-1 else LIGHT)
    add_rect(s, x, y0 + chart_h - h, bar_w, h, color)
    # Value label
    add_text(s, x - Inches(0.3), y0 + chart_h - h - Inches(0.45),
             bar_w + Inches(0.6), Inches(0.4),
             f"{m:.1f}tr", size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    # District name
    add_text(s, x - Inches(0.4), y0 + chart_h + Inches(0.1),
             bar_w + Inches(0.8), Inches(0.6),
             dist, size=10, color=TEXT, align=PP_ALIGN.CENTER)

# Nhận xét
add_text(s, Inches(0.5), Inches(6.4), Inches(12), Inches(0.5),
         "TRẢ LỜI:", size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.6),
         "Quận Bình Thạnh (71.4 tr/m²) cao nhất - gần trung tâm, view sông. Bình Tân (42.8 tr/m²) thấp nhất.",
         size=14, color=TEXT)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 21: Trả lời C2, C3 - Quan hệ
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "C2/C3: DIỆN TÍCH & HƯỚNG NHÀ ẢNH HƯỞNG?", page=21, total=total_slides)
# Left: area vs price scatter
add_text(s, Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
         "Diện tích ↔ Tổng giá", size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_image(s, FIG_DIR / "fig04_area_vs_price.png",
          Inches(0.5), Inches(1.7), w=Inches(6))
# Right: bedroom distribution
add_text(s, Inches(6.8), Inches(1.3), Inches(6), Inches(0.4),
         "Phân bố số phòng ngủ", size=15, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_image(s, FIG_DIR / "fig10_bedrooms_count.png",
          Inches(6.8), Inches(1.7), w=Inches(6))
# Trả lời
add_text(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.4),
         "TRẢ LỜI:", size=14, bold=True, color=ACCENT)
add_bullet_box(s, Inches(0.7), Inches(6.0), Inches(12), Inches(1.4), [
    "Diện tích và tổng giá có quan hệ gần tuyến tính (corr=0.70); nhưng giá/m² thì không phụ thuộc nhiều vào diện tích.",
    "2PN chiếm đa số (67%, 1189 tin) - cơ cấu căn hộ TP.HCM; 1PN và 3PN chiếm ~16% mỗi loại.",
    "Hướng nhà, ban công, nội thất, pháp lý có ảnh hưởng YẾU đến giá (corr<0.3); vị trí (district) và dự án quan trọng hơn.",
], size=12)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 22: Trả lời C4 - Kết quả 4 mô hình (Table 4)
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "C4: HIỆU SUẤT MÔ HÌNH DỰ ĐOÁN", page=22, total=total_slides)
# Big table
rows = [
    ("Dummy (Baseline)", "18.8M", "34.1M", "−0.048", "18.5M", "37.1M", "−0.041"),
    ("Linear Regression", "14.1M", "28.1M", "0.294", "14.9M", "33.7M", "0.143"),
    ("Random Forest ★", "13.0M", "26.8M", "0.357", "13.9M", "33.2M", "0.169"),
    ("Gradient Boosting", "13.4M", "26.7M", "0.366", "14.4M", "33.3M", "0.162"),
]
headers = ["Mô hình", "CV MAE", "CV RMSE", "CV R²", "Test MAE", "Test RMSE", "Test R²"]
table = s.shapes.add_table(5, 7, Inches(0.5), Inches(1.3),
                             Inches(12.3), Inches(2.8)).table
col_w = [Inches(2.5), Inches(1.5), Inches(1.7), Inches(1.6), Inches(1.7), Inches(1.7), Inches(1.6)]
for i, w in enumerate(col_w):
    table.columns[i].width = w
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True
            r.font.color.rgb = WHITE
            r.font.size = Pt(13)
            r.font.name = "Calibri"
for ri, row in enumerate(rows):
    bg = LIGHT if ri % 2 == 0 else WHITE
    is_best = "★" in row[0]
    for ci, val in enumerate(row):
        cell = table.cell(ri+1, ci)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN if is_best and ci > 0 else bg
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(12)
                r.font.color.rgb = WHITE if (is_best and ci > 0) else TEXT
                r.font.bold = is_best or ci == 0
                r.font.name = "Calibri"

# Nhận xét
add_text(s, Inches(0.5), Inches(4.4), Inches(12), Inches(0.4),
         "Chỉ số đo lường: MAE (Mean Absolute Error) là sai số tuyệt đối trung bình; RMSE là căn bậc hai sai số bình phương; R² là hệ số xác định (1.0 = hoàn hảo).",
         size=11, italic=True, color=SUBTLE)
add_text(s, Inches(0.5), Inches(4.9), Inches(12), Inches(0.5),
         "TRẢ LỜI:", size=14, bold=True, color=ACCENT)
add_bullet_box(s, Inches(0.7), Inches(5.3), Inches(12), Inches(2), [
    "Random Forest (mô hình tốt nhất): Test R² = 0.169, MAE ≈ 13.9 triệu VND/m² (~26% sai số so với median 53 tr/m²).",
    "Dummy R² âm (-0.041) đúng kỳ vọng (median không dự đoán được). Linear R² = 0.143 giải thích ~14% phương sai.",
    "RF > Linear (+0.026 R²) chứng minh có quan hệ phi tuyến. GBR CV R² = 0.366 nhỉnh hơn nhưng gap CV-Test = 0.204 > RF (0.188) → RF generalize tốt hơn.",
    "R² thấp (0.17) là thực tế của BĐS: giá phụ thuộc nhiều yếu tố phi số (view, tầng, nội thất chi tiết) không có trong data.",
], size=12)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 23: Trả lời C5 - Phân cụm (Table 6, Table 7, Hình 6)
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "C5: PHÂN KHÚC THỊ TRƯỜNG (K-MEANS)", page=23, total=total_slides)
# Left: silhouette table
add_text(s, Inches(0.5), Inches(1.3), Inches(5.5), Inches(0.4),
         "Chọn K tự động (Silhouette Score)", size=14, bold=True, color=DARK)
table_s = s.shapes.add_table(5, 3, Inches(0.5), Inches(1.8),
                              Inches(5.5), Inches(2.2)).table
table_s.columns[0].width = Inches(1.0)
table_s.columns[1].width = Inches(2.5)
table_s.columns[2].width = Inches(2.0)
for i, h in enumerate(["K", "Silhouette", "Chọn"]):
    cell = table_s.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True
            r.font.color.rgb = WHITE
            r.font.size = Pt(12)
            r.font.name = "Calibri"
ks = [("3", "0.079", "—"), ("4", "0.083", "✓ tốt nhất"), ("5", "0.027", "—"), ("6", "−0.061", "—")]
for ri, row in enumerate(ks):
    is_best = "✓" in row[2]
    for ci, val in enumerate(row):
        cell = table_s.cell(ri+1, ci)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN if is_best else (LIGHT if ri % 2 == 0 else WHITE)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(12)
                r.font.color.rgb = WHITE if is_best else TEXT
                r.font.bold = is_best
                r.font.name = "Calibri"

# Right: cluster distribution table
add_text(s, Inches(6.3), Inches(1.3), Inches(6.5), Inches(0.4),
         "Phân bố cụm (K=4, n=1773)", size=14, bold=True, color=DARK)
table_c = s.shapes.add_table(5, 4, Inches(6.3), Inches(1.8),
                              Inches(6.5), Inches(2.2)).table
for i, w in enumerate([Inches(0.9), Inches(1.0), Inches(1.4), Inches(3.2)]):
    table_c.columns[i].width = w
for i, h in enumerate(["Cụm", "Số tin", "Tỷ lệ", "Đặc điểm"]):
    cell = table_c.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True
            r.font.color.rgb = WHITE
            r.font.size = Pt(12)
            r.font.name = "Calibri"
clusters = [
    ("0", "165", "9.3%", "Hướng Đông, balcony Đông"),
    ("1", "1396", "78.7%", "Đa số - căn 2PN đặc trưng"),
    ("2", "47", "2.6%", "Căn nhỏ ở Thủ Đức (~55m²)"),
    ("3", "165", "9.3%", "Hướng Tây/TB, balcony Tây"),
]
for ri, row in enumerate(clusters):
    for ci, val in enumerate(row):
        cell = table_c.cell(ri+1, ci)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT if ri % 2 == 0 else WHITE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if ci < 3 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(11)
                r.font.color.rgb = TEXT
                r.font.bold = (ci == 0)
                r.font.name = "Calibri"

# Bottom: silhouette figure + trả lời
add_image(s, FIG_DIR / "fig11_silhouette.png",
          Inches(0.5), Inches(4.2), w=Inches(5))
add_text(s, Inches(0.5), Inches(4.0), Inches(5), Inches(0.3),
         "Hình 3: Silhouette score theo K", size=11, italic=True,
         color=SUBTLE, align=PP_ALIGN.CENTER)

add_text(s, Inches(6.3), Inches(4.2), Inches(6.5), Inches(0.4),
         "TRẢ LỜI:", size=14, bold=True, color=ACCENT)
add_bullet_box(s, Inches(6.3), Inches(4.6), Inches(6.5), Inches(2.7), [
    "Thị trường phân thành K=4 phân khúc.",
    "Cluster 1 chiếm đa số 78.7% (1396 tin).",
    "Cluster 2 chỉ 47 tin (2.6%) tập trung ở Thủ Đức, căn nhỏ ~55m².",
    "Các cụm khác biệt theo direction_code (hướng nhà), KHÔNG theo giá (median 50–54M gần như nhau).",
    "Silhouette thấp (0.083) → cụm chồng lấn, không tách rõ phân khúc giá.",
], size=12)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 24: Trả lời C6 - Top 5 Recommender (Table 8, Hình 7)
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "C6: HỆ GỢI Ý TOP-5 CĂN HỘ", page=24, total=total_slides)
# Left: profiles
add_text(s, Inches(0.5), Inches(1.3), Inches(6.5), Inches(0.4),
         "3 hồ sơ nhu cầu mẫu", size=14, bold=True, color=DARK)
table_p = s.shapes.add_table(4, 4, Inches(0.5), Inches(1.8),
                              Inches(6.5), Inches(2.4)).table
for i, w in enumerate([Inches(2.0), Inches(1.2), Inches(2.5), Inches(0.8)]):
    table_p.columns[i].width = w
for i, h in enumerate(["Hồ sơ", "Ngân sách", "Quận ưu tiên", "Kết quả"]):
    cell = table_p.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True
            r.font.color.rgb = WHITE
            r.font.size = Pt(12)
            r.font.name = "Calibri"
profiles = [
    ("Gia đình trẻ (2PN, 65m²)", "3 tỷ", "Thủ Đức, B.Thạnh", "Top-5 ✓"),
    ("Nhà đầu tư (2PN, 70m²)", "5 tỷ", "Quận 7, B.Tân", "Top-5 ✓"),
    ("Người mua cao cấp (3PN, 85m²)", "7 tỷ", "Thủ Đức, Q7", "Top-5 ✓"),
]
for ri, row in enumerate(profiles):
    bg = GREEN if "✓" in row[3] else LIGHT
    for ci, val in enumerate(row):
        cell = table_p.cell(ri+1, ci)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(11)
                r.font.color.rgb = WHITE if "✓" in row[3] and ci == 3 else TEXT
                r.font.bold = (ci == 0) or ("✓" in row[3] and ci == 3)
                r.font.name = "Calibri"

# Right: amenity
add_text(s, Inches(7.3), Inches(1.3), Inches(5.5), Inches(0.4),
         "Điểm tiện ích theo quận", size=14, bold=True, color=DARK)
add_image(s, FIG_DIR / "fig09_amenity_by_district.png",
          Inches(7.3), Inches(1.8), w=Inches(5.5))

# Trả lời
add_text(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.4),
         "TRẢ LỜI:", size=14, bold=True, color=ACCENT)
add_bullet_box(s, Inches(0.7), Inches(4.8), Inches(12.3), Inches(2.4), [
    "Cả 3 hồ sơ đều có kết quả top-5 ✓ (recommender hoạt động trơn tru).",
    "Profile 1 (3 tỷ): 5/5 tin thuộc Thủ Đức, giá 2.7–3.2 tỷ (sát budget).",
    "Profile 2 (5 tỷ): 5/5 tin Quận 7, cluster 1, giá 4.8–5.1 tỷ.",
    "Profile 3 (7 tỷ): cluster 2 chỉ 47 tin Thủ Đức (giá thấp 1.6–4.6 tỷ) → 0 khớp → chuyển cluster 1.",
    "Công thức: score = price_score + area_score + cluster_bonus(0.3) + amenity_bonus(0.2×amenity/max).",
], size=13)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 25: Phân tích 10 worst (Table 5)
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "TOP 10 TRƯỜNG HỢP DỰ ĐOÁN SAI LỚN NHẤT", page=25, total=total_slides)
rows_w = [
    ("1", "177641859", "Quận 12", "26.0", "1", "480.8", "42.2", "91.2%"),
    ("2", "177622744", "Q.Bình Thạnh", "60.0", "3", "258.3", "57.5", "77.8%"),
    ("3", "177494566", "Q.Bình Thạnh", "122.4", "3", "228.8", "72.4", "68.4%"),
    ("4", "177850739", "TP.Thủ Đức", "140.0", "3", "51.4", "127.3", "147.4%"),
    ("5", "177757531", "TP.Thủ Đức", "74.0", "2", "162.2", "58.9", "63.7%"),
    ("6", "177715817", "Quận 7", "82.0", "2", "146.3", "54.8", "62.5%"),
    ("7", "177770338", "Q.Bình Thạnh", "65.0", "2", "121.5", "55.9", "54.0%"),
    ("8", "177866544", "TP.Thủ Đức", "51.0", "1", "117.6", "52.4", "55.4%"),
    ("9", "177766301", "TP.Thủ Đức", "71.0", "2", "133.8", "70.2", "47.5%"),
    ("10", "175165640", "TP.Thủ Đức", "64.0", "1", "215.6", "98.5", "54.3%"),
]
headers_w = ["#", "listing_id", "Quận", "DT(m²)", "PN", "Giá thực(tr)", "Dự đoán(tr)", "%Sai"]
table_w = s.shapes.add_table(11, 8, Inches(0.5), Inches(1.3),
                              Inches(12.3), Inches(3.6)).table
col_w_w = [Inches(0.5), Inches(1.7), Inches(1.7), Inches(1.0), Inches(0.6),
           Inches(2.0), Inches(2.0), Inches(2.8)]
for i, w in enumerate(col_w_w):
    table_w.columns[i].width = w
for i, h in enumerate(headers_w):
    cell = table_w.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True
            r.font.color.rgb = WHITE
            r.font.size = Pt(11)
            r.font.name = "Calibri"
for ri, row in enumerate(rows_w):
    bg = LIGHT if ri % 2 == 0 else WHITE
    pct = float(row[7].rstrip("%"))
    high_pct = pct > 60
    for ci, val in enumerate(row):
        cell = table_w.cell(ri+1, ci)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if high_pct and ci == 7 else bg
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(10)
                r.font.color.rgb = WHITE if high_pct and ci == 7 else TEXT
                r.font.bold = (ci == 0) or high_pct and ci == 7
                r.font.name = "Calibri"

# Phân tích
add_text(s, Inches(0.5), Inches(5.1), Inches(12), Inches(0.5),
         "PHÂN TÍCH:", size=14, bold=True, color=ACCENT)
add_bullet_box(s, Inches(0.7), Inches(5.5), Inches(12), Inches(2), [
    "5/10 căn giá cực cao (>150 tr/m²): penthouse/căn góc view đẹp - model underprice 30-60% do không có feature view/tầng.",
    "1 trường hợp (STT 4): căn Duplex 140m² Thủ Đức, actual 51.4 tr/m² bất thường thấp → có thể data entry error trong tin đăng gốc.",
    "5/10 thiếu ≥3 features numeric → median imputation gây sai lệch cho 50% top-10 worst.",
], size=13)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 26: SECTION - KẾT LUẬN
# ─────────────────────────────────────────────────────────────────────────
s = slide_section("KẾT LUẬN & HƯỚNG PHÁT TRIỂN")

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 27: Kết luận
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "KẾT LUẬN - TỔNG KẾT KẾT QUẢ", page=27, total=total_slides)
add_text(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
         "Đồ án đã hoàn thành trọn vẹn quy trình CRISP-DM, giải quyết các mục tiêu:",
         size=14, italic=True, color=SUBTLE)
add_bullet_box(s, Inches(0.7), Inches(1.9), Inches(12), Inches(5.5), [
    ("Pipeline tự động: API crawling + 5 bước làm sạch → 1.773 dòng sạch từ 1.799 tin.", 15),
    ("EDA: Bình Thạnh cao nhất (71.4 tr/m²), Bình Tân thấp nhất (42.8 tr/m²); 2PN chiếm 67%.", 15),
    ("Random Forest tối ưu: Test R² = 0.169, MAE = 13.9 tr/m² (~26% sai số so với median 53M).", 15),
    ("K-Means K=4 (Silhouette = 0.083): 4 cụm khác biệt theo direction_code, KHÔNG theo phân khúc giá.", 15),
    ("Hệ gợi ý hybrid: 3/3 hồ sơ đều có top-5 đúng ràng buộc (budget ±20%, quận, phòng ±1).", 15),
    ("45/45 unit test PASS, 4 Notebook Jupyter chạy end-to-end, báo cáo Markdown + slide + AI log đầy đủ.", 15),
], size=15)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 28: Hạn chế
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "HẠN CHẾ CỦA NGHIÊN CỨU", page=28, total=total_slides)
add_bullet_box(s, Inches(0.7), Inches(1.4), Inches(12), Inches(5.5), [
    ("Không gian & thời gian: 1.773 quan sát snapshot T7-T8/2026, 6 quận - chưa đại diện 24 quận TP.HCM, không phân tích xu hướng.", 16),
    ("Biến số ẩn: R² = 0.17 và Silhouette = 0.083 khiêm tốn - giá BĐS bị chi phối bởi nhiều yếu tố vô hình (view sông, tầng cao, chất lượng nội thất) không có trong data.", 16),
    ("Chất lượng dữ liệu: Missing cao (direction 79%, balcony 73%) → median imputation có thể bias. Có 1 case Duplex data entry error (147% sai số) khó phát hiện.", 16),
    ("Cluster mất cân đối: cluster 1 = 78.7%, cluster 2 = 2.6% → bonus cluster không hiệu quả cho cluster nhỏ.", 16),
    ("Hệ gợi ý đơn giản: rule-based scoring, chưa có collaborative filtering do thiếu dữ liệu hành vi người dùng.", 16),
    ("Không khai thác văn bản: title/description có thể chứa thông tin view/tầng nhưng chưa được text mining.", 16),
], size=16)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 29: Hướng phát triển
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "KIẾN NGHỊ & HƯỚNG PHÁT TRIỂN", page=29, total=total_slides)
add_bullet_box(s, Inches(0.7), Inches(1.4), Inches(12), Inches(5.5), [
    ("Mở rộng dữ liệu: Tự động hóa crawling đa nguồn (batdongsan, mogi, alonhadat), hướng tới ≥10.000 quan sát bao phủ 24 quận TP.HCM.", 16),
    ("Feature engineering: Tích hợp GIS từ OpenStreetMap tính khoảng cách đến trung tâm, bổ sung bedroom_density, is_high_end_district, log(area_m²).", 16),
    ("Text mining: Ứng dụng LLM/TF-IDF trích xuất thuộc tính ẩn (tầng, view, nội thất chi tiết) từ title và description.", 16),
    ("Mô hình tiên tiến: Thử XGBoost, LightGBM với Optuna hyperparameter tuning. Thêm feature missing_count để model biết uncertainty.", 16),
    ("Hệ gợi ý nâng cấp: Mở rộng 'quận tương đương' (cùng tier giá), thêm trọng số tuỳ biến, thử collaborative filtering khi có đủ dữ liệu.", 16),
    ("Triển khai: REST API (Flask/FastAPI) + dashboard trực quan (Streamlit) phục vụ demo cho người dùng cuối.", 16),
], size=16)

# ─────────────────────────────────────────────────────────────────────────
# SLIDE 30: CÂU HỎI & TÀI LIỆU THAM KHẢO
# ─────────────────────────────────────────────────────────────────────────
s = slide_blank()
slide_header(s, "TÀI LIỆU THAM KHẢO & LỜI CẢM ƠN", page=30, total=total_slides)
add_text(s, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
         "TÀI LIỆU THAM KHẢO", size=18, bold=True, color=DARK)
refs_short = [
    "[1] Pedregosa, F. et al. (2011). Scikit-learn: Machine Learning in Python. JMLR 12:2825–2830.",
    "[4] Breiman, L. (2001). Random Forests. Machine Learning 45(1):5–32.",
    "[5] Friedman, J.H. (2001). Greedy Function Approximation: Gradient Boosting. Annals of Statistics 29(5).",
    "[6] Rousseeuw, P.J. (1987). Silhouettes: A Graphical Aid to Cluster Analysis. J. Computational & Applied Math 20:53–65.",
    "[7] Ricci, F., Rokach, L., Shapira, B. (2015). Recommender Systems Handbook (2nd ed.). Springer.",
    "[8] James, G. et al. (2013). An Introduction to Statistical Learning. Springer.",
    "[9] Chen, T., Guestrin, C. (2016). XGBoost: A Scalable Tree Boosting System. KDD '16:785–794.",
    "[11] Đề tài Chuyên đề cuối kỳ - LTCKHDL (file đề bài).",
    "[12] Dataset real_estate_apartment.xlsx - snapshot chotot.com (T7-T8/2026).",
]
add_bullet_box(s, Inches(0.7), Inches(1.8), Inches(6.2), Inches(5.5), refs_short, size=10)

# Right: Lời cảm ơn
add_rect(s, Inches(7.2), Inches(1.3), Inches(5.6), Inches(5.5), DARK)
add_text(s, Inches(7.2), Inches(1.5), Inches(5.6), Inches(0.5),
         "LỜI CẢM ƠN", size=22, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, Inches(7.4), Inches(2.3), Inches(5.2), Inches(2.5),
         "Em xin chân thành cảm ơn PGS.TS Nguyễn Duy Hàm đã hướng dẫn tận tình trong suốt quá trình thực hiện đồ án.\n\nEm cũng xin cảm ơn các bạn học viên cao học lớp 33CNTT21-PH đã hỗ trợ và trao đổi kinh nghiệm.",
         size=14, color=WHITE, italic=True, align=PP_ALIGN.JUSTIFY)
add_text(s, Inches(7.2), Inches(6.0), Inches(5.6), Inches(0.6),
         "Xin trân trọng cảm ơn!",
         size=20, bold=True, color=ACCENT,
         align=PP_ALIGN.CENTER)

# ── Lưu ──────────────────────────────────────────────────────────────────
prs.save(str(OUT))
print(f"✅ Đã lưu: {OUT}")
print(f"   Kích thước: {OUT.stat().st_size / 1024:.1f} KB")
print(f"   Số slides: {len(prs.slides)}")